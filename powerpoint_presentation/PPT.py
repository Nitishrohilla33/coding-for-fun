from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- Create widescreen (16:9) presentation ----
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = 13.333   # slide width in inches
SH = 7.5      # slide height in inches
MARGIN = 0.8  # consistent page margin

# ---- Color Palette ----
DEEP_BLUE = RGBColor(0, 51, 102)
TEAL      = RGBColor(0, 102, 153)
DARK_GREY = RGBColor(60, 60, 60)
BLACK     = RGBColor(0, 0, 0)

LOGO = "/Users/mac/Desktop/iit-tirupati-india-logo-png_seeklogo-464788.png"


# -------- Helper: clear default placeholders --------
def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


# -------- Helper: format text --------
def format_text(text_frame, title=False):
    for para in text_frame.paragraphs:
        # Paragraph-level properties - set OUTSIDE the run loop
        para.alignment = PP_ALIGN.CENTER if title else PP_ALIGN.LEFT
        para.space_after = Pt(8) if title else Pt(12)
        para.line_spacing = 1.4 if title else 1.6

        for run in para.runs:
            run.font.name = "Times New Roman"
            if title:
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = DEEP_BLUE
            else:
                run.font.size = Pt(20)
                run.font.color.rgb = DARK_GREY


# -------- Content Slide Creator --------
def add_bullet_slide(title_text, points, image_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    clear_slide(slide)

    # Top accent stripe (full width)
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(SW), Inches(0.08)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL
    top_bar.line.fill.background()

    # Left accent bar
    side_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN - 0.3), Inches(1.6),
        Inches(0.12), Inches(4.2)
    )
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = TEAL
    side_bar.line.fill.background()

    # Logo (top-right)
    if image_path:
        slide.shapes.add_picture(
            image_path,
            Inches(SW - MARGIN - 0.8), Inches(0.15),
            width=Inches(0.8)
        )

    # Title - centered across full content width
    title_w = SW - 2 * MARGIN
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(0.35),
        Inches(title_w), Inches(1)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = title_text
    format_text(tf, title=True)

    # Divider under title
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN), Inches(1.4),
        Inches(title_w), Inches(0.025)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = TEAL
    divider.line.fill.background()

    # Bullet points
    content_box = slide.shapes.add_textbox(
        Inches(MARGIN + 0.2), Inches(1.8),
        Inches(title_w - 0.4), Inches(5)
    )
    ctf = content_box.text_frame
    ctf.word_wrap = True
    ctf.clear()

    for i, point in enumerate(points):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = f"\u2022   {point}"
        p.level = 0

    format_text(ctf)

    return slide


# ================================================================
#                         TITLE SLIDE
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[5])
clear_slide(slide)

# Full-width top bar
bar_top = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0),
    Inches(SW), Inches(0.22)
)
bar_top.fill.solid()
bar_top.fill.fore_color.rgb = TEAL
bar_top.line.fill.background()

# Full-width bottom bar
bar_bot = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(SH - 0.22),
    Inches(SW), Inches(0.22)
)
bar_bot.fill.solid()
bar_bot.fill.fore_color.rgb = TEAL
bar_bot.line.fill.background()

# Logo (top-left)
slide.shapes.add_picture(
    LOGO,
    Inches(MARGIN), Inches(0.45),
    width=Inches(1.2)
)

# Date (top-right) -- positioned WITHIN slide bounds
date_box = slide.shapes.add_textbox(
    Inches(SW - MARGIN - 3), Inches(0.5),
    Inches(3), Inches(0.5)
)
date_tf = date_box.text_frame
date_tf.text = "April 2026"
date_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
for run in date_tf.paragraphs[0].runs:
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_GREY

# Title (centered, each line as its own paragraph for proper centering)
title_w = SW - 2 * MARGIN
title_box = slide.shapes.add_textbox(
    Inches(MARGIN), Inches(2.0),
    Inches(title_w), Inches(2.5)
)
title_tf = title_box.text_frame
title_tf.word_wrap = True

lines = [
    "Light\u2013Material Interactions Using Laser and Flash",
    "Sources for Energy Conversion and Storage",
    "Applications"
]
for i, line in enumerate(lines):
    p = title_tf.paragraphs[0] if i == 0 else title_tf.add_paragraph()
    p.text = line
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)
    p.line_spacing = 1.4
    for run in p.runs:
        run.font.name = "Times New Roman"
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = DEEP_BLUE

# Accent line above name
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(MARGIN), Inches(5.0),
    Inches(3.5), Inches(0.035)
)
accent.fill.solid()
accent.fill.fore_color.rgb = TEAL
accent.line.fill.background()

# Name
name_box = slide.shapes.add_textbox(
    Inches(MARGIN), Inches(5.2),
    Inches(6), Inches(0.5)
)
name_tf = name_box.text_frame
name_tf.text = "Nitish Rohilla (PH26D002)"
name_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
for run in name_tf.paragraphs[0].runs:
    run.font.name = "Times New Roman"
    run.font.size = Pt(20)
    run.font.color.rgb = BLACK
    run.font.bold = True

# Department
dept_box = slide.shapes.add_textbox(
    Inches(MARGIN), Inches(5.75),
    Inches(6), Inches(0.5)
)
dept_tf = dept_box.text_frame
dept_tf.text = "Department of Physics"
dept_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
for run in dept_tf.paragraphs[0].runs:
    run.font.name = "Times New Roman"
    run.font.size = Pt(16)
    run.font.color.rgb = DARK_GREY


# ================================================================
#                        CONTENT SLIDES
# ================================================================

add_bullet_slide("Motivation", [
    "Increasing global energy demand",
    "Limitations of traditional processing",
    "Need for faster material techniques",
    "Development of advanced energy devices"
], image_path=LOGO)

add_bullet_slide("Objectives", [
    "Understand light\u2013material interactions",
    "Study laser and flash processing",
    "Analyze photothermal and photochemical effects",
    "Explore energy applications"
], image_path=LOGO)

add_bullet_slide("Methodology", [
    "Laser and flash light processing",
    "Control of wavelength and intensity",
    "Photothermal processes",
    "Photochemical modifications"
], image_path=LOGO)

add_bullet_slide("Key Results", [
    "Faster processing than furnace methods",
    "Improved material properties",
    "Enhanced device performance",
    "Scalable flash processing"
], image_path=LOGO)

add_bullet_slide("Limitations", [
    "Complex interaction mechanisms",
    "Difficult parameter control",
    "Limited scalability in some methods"
], image_path=LOGO)

add_bullet_slide("Thank You", [
    "Thank you for your attention"
], image_path=LOGO)


# ================================================================
#                            SAVE
# ================================================================
prs.save("/Users/mac/Desktop/LMI_Presentation.pptx")
print("PPT created successfully on Desktop!")
